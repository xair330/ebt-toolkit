"""
PDF → PPT 转换脚本
每页 PDF 渲染为高清图片后嵌入幻灯片，比例自动适配
"""
import pathlib, io, sys
import pymupdf                         # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ── 路径 ────────────────────────────────────────────────────────
PDF_PATH = pathlib.Path(r"c:\工作文件\数据管理\检查员\2026一季度 EBT胜任力数据深度分析报告.pdf")
PPT_PATH = PDF_PATH.with_suffix(".pptx")

# ── 渲染参数 ─────────────────────────────────────────────────────
DPI   = 180          # 分辨率（越高越清晰但文件越大）
SCALE = DPI / 72     # PDF 基准分辨率是 72 dpi

print(f"[PDF→PPT] 读取: {PDF_PATH.name}")
doc = pymupdf.open(str(PDF_PATH))
total = len(doc)
print(f"[PDF→PPT] 共 {total} 页，开始转换...")

# 用第一页的尺寸确定幻灯片宽高
first_page = doc[0]
w_pt = first_page.rect.width   # PDF 点（1pt = 1/72 inch）
h_pt = first_page.rect.height

# PPT 幻灯片尺寸（EMU = English Metric Unit，1英寸=914400 EMU）
slide_w = int(w_pt / 72 * 914400)
slide_h = int(h_pt / 72 * 914400)

prs = Presentation()
prs.slide_width  = slide_w
prs.slide_height = slide_h

blank_layout = prs.slide_layouts[6]  # 完全空白布局

for i, page in enumerate(doc):
    sys.stdout.write(f"\r  渲染第 {i+1}/{total} 页...")
    sys.stdout.flush()

    # 渲染为 RGB 图像
    mat = pymupdf.Matrix(SCALE, SCALE)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img_bytes = pix.tobytes("png")

    # 创建幻灯片
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        io.BytesIO(img_bytes),
        left=0, top=0,
        width=slide_w, height=slide_h
    )

doc.close()
sys.stdout.write("\n")

prs.save(str(PPT_PATH))
print(f"[PDF→PPT] 完成！输出文件：")
print(f"  {PPT_PATH}")
print(f"  大小约 {PPT_PATH.stat().st_size / 1024 / 1024:.1f} MB")
